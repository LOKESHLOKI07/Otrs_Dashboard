from datetime import datetime
import datetime
from pptx.dml.color import ColorFormat, RGBColor


import mysql.connector
from cssutils.css.cssvalue import RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd

pd.set_option('display.max_columns', 10)
# Connect to the MySQL database
conn = mysql.connector.connect(
    host='otrs.futurenet.in',
    port=3306,
    user='readuser2',
    password='6FbUDa5VM',
    database='otrs5'
)
cursor = conn.cursor()
# get current date
today = datetime.date.today()

# get last month
last_month = today - datetime.timedelta(days=30)

# get two months ago
two_months_ago = today - datetime.timedelta(days=60)

# get three months ago
three_months_ago = today - datetime.timedelta(days=90)

# get month names
last_month_name = last_month.strftime('%B')
two_months_ago_name = two_months_ago.strftime('%B')
three_months_ago_name = three_months_ago.strftime('%B')

# Create the PowerPoint presentation
prs = Presentation()

# add a new slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

logo_path = '/home/ubuntu/Downloads/futurenet-logo.png'
logo = slide.shapes.add_picture(logo_path, left=Inches(0.5), top=Inches(1.5), height=Inches(1))

# add a header to the slide
header = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
header.text = 'CUSTOMER NAME Monthly Review Meeting'
header.text_frame.paragraphs[0].font.size = Inches(0.4)

# add a body to the slide
body = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(5))
body.text = 'Remote Infrastructure Management Services (RIMS)'
body.text_frame.paragraphs[0].font.size = Inches(0.3)

# add a footer to the slide
# today = datetime.now().strftime("%d/%m/%Y")
today = datetime.date.today()
today_str = today.strftime('%Y-%m-%d')
print(today)
footer = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(9), Inches(0.5))
footer.text = 'Date:' + ' ' + today_str

# Set up the title slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
body = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(5))
body.text = 'Ticket Statistics'
body.text_frame.paragraphs[0].font.size = Inches(0.5)
# title.text = "Ticket Statistics"


# Execute a new SQL query to retrieve the desired data
new_query = """ SELECT  tt.name,
  SUM(CASE WHEN ts.id IN (1,4,6,7,8,11,13, 15, 16, 14, 17,18) and t.create_time BETWEEN DATE_SUB(NOW(), INTERVAL 1 day) AND NOW() THEN 1 ELSE 0 END) AS pending,
SUM(CASE WHEN ts.id IN (2, 3,10,12) and t.create_time BETWEEN DATE_SUB(NOW(), INTERVAL 1 day) AND NOW() THEN 1 ELSE 0 END) AS closed,
SUM(CASE WHEN ts.id IN (2, 3,10,12) and t.change_time BETWEEN DATE_SUB(NOW(), INTERVAL 1 day) AND NOW() THEN 1 ELSE 0 END) AS 'overallClosed'
FROM 
  ticket AS t 
left JOIN ticket_state ts ON t.ticket_state_id = ts.id
left JOIN users u ON t.user_id = u.id
left JOIN ticket_priority tp ON t.ticket_priority_id = tp.id
left JOIN ticket_type tt ON t.type_id = tt.id
left JOIN queue q ON t.queue_id = q.id
WHERE 

 ts.name NOT IN ('merged') 
  AND q.name NOT IN ('SALES', 'PRESALES', 'ODOOHELPDESK') 
GROUP BY tt.name 
HAVING (pending > 0 OR closed > 0 OR overallClosed > 0) """
# Your new SQL query here
cursor.execute(new_query)
new_result = cursor.fetchall()

# Set title for slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
body = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(5))
body.text = 'Overview Of Ticket'
body.text_frame.paragraphs[0].font.size = Inches(0.5)

# Create a table shape on the new slide
rows = len(new_result) + 1  # Include header row
cols = len(new_result[0])
left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(1)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Populate the header row of the table
header_row = table.rows[0].cells
header_row[0].text = "CATEGORY"
header_row[1].text = "Registered"
header_row[2].text = "CLOSED"
header_row[3].text = "WORK IN PROCESS"
# ...

# Populate the remaining rows of the table with the retrieved data
for i, row in enumerate(new_result):
    table_row = table.rows[i + 1].cells
    table_row[0].text = str(row[0])
    table_row[1].text = str(row[1] + row[2])
    table_row[2].text = str(row[2])
    table_row[3].text = str(row[1])

# end of 2nd slide
# 3rd slide


# Execute a new SQL query to retrieve the desired data
catergory_query = """ SELECT tt.name,
SUM(CASE WHEN tt.id IN (2) THEN 1 ELSE 0 END) AS 'INCIDENT',
SUM(CASE WHEN tt.id IN (14) THEN 1 ELSE 0 END) AS 'PROACTIVE',
SUM(CASE WHEN tt.id IN (4) THEN 1 ELSE 0 END) AS 'SERVICE'
#SUM(CASE WHEN ts.id IN (1,4,6,7,8,11,13, 15, 16, 14, 17,18)  THEN 1 ELSE 0 END) AS 'Pending'

FROM ticket t
JOIN ticket_state ts ON t.ticket_state_id = ts.id
JOIN users u ON t.user_id = u.id
JOIN ticket_type tt ON t.type_id = tt.id
JOIN queue q ON t.queue_id = q.id
WHERE  tt.name NOT IN ('junk') AND ts.name NOT IN ('merged')
AND q.name NOT IN ('SALES','PRESALES','ODOOHELPDESK')
AND t.create_time BETWEEN DATE_SUB(NOW(), INTERVAL 1 day) AND NOW()
GROUP BY tt.name """
# Your new SQL query here
cursor.execute(catergory_query)
catergory_result = cursor.fetchall()
max_row = None
max_value = 0

for row in catergory_result:
    for value in row[1:]:
        if value > max_value:
            max_value = value
            max_row = row

print("Row with maximum value:", max_row)
print(max_row[0])

# Set title for slide 2
# Set title for slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
body = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(3))
body.text = 'Ticket Category-Technology Wise'
body.text_frame.paragraphs[0].font.size = Inches(0.5)

# add a footer to the slide
# today = datetime.now().strftime("%d/%m/%Y")
footer = slide.shapes.add_textbox(Inches(0.1), Inches(6.5), Inches(10), Inches(0.5))
footer.text = (last_month_name + ' -2023  ' + max_row[0] + ' ALERT RELATED – ' + str(max_value) + ' CALLS').upper()
# table_font_size = Pt(15)
footer.text_frame.paragraphs[0].font.size = Pt(18)

# Create a table shape on the new slide
rows = len(catergory_result) + 1  # Include header row
cols = len(catergory_result[0]) + 1

left, top, width, height = Inches(1), Inches(1.5), Inches(7), Inches(0.5)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Populate the header row of the table
header_row = table.rows[0].cells

header_row[0].text = "Category"
header_row[1].text = "Incident"
header_row[2].text = "Proactive Notification"
header_row[3].text = "Service Request"
header_row[4].text = "Total"
# ...

# Populate the remaining rows of the table with the retrieved data
for i, row in enumerate(catergory_result):
    table_row = table.rows[i + 1].cells
    table_row[0].text = str(row[0])
    table_row[1].text = str(row[1])
    table_row[2].text = str(row[2])
    table_row[3].text = str(row[3])
    table_row[4].text = str(row[1] + row[2] + row[3])
 # to set size
table_font_size = Pt(15)
for cell in table.iter_cells():
    cell.text_frame.paragraphs[0].font.size = table_font_size
# end of 3rd slide

# 4th slide

# Connect to database and fetch data

tp_query = '''SELECT  s.name as 'INCIDENT & NOTIFICATION',
SUM(CASE WHEN ts.id IN (1, 4, 6, 7, 8, 11, 13, 15, 16, 14, 17, 18, 2, 3, 10, 12) THEN 1 ELSE 0 END) AS 'TOTAL',
SUM(CASE WHEN ts.id IN (2, 3, 10, 12) THEN 1 ELSE 0 END) AS 'CLOSED',
SUM(CASE WHEN ts.id IN (1, 4, 6, 7, 8, 11, 13, 15, 16, 14, 17, 18) THEN 1 ELSE 0 END) AS 'Work In Process'
FROM ticket t
JOIN ticket_state ts ON t.ticket_state_id = ts.id
JOIN users u ON t.user_id = u.id
JOIN ticket_type tt ON t.type_id = tt.id
JOIN queue q ON t.queue_id = q.id
left JOIN sla s ON s.id = t.sla_id
WHERE  tt.name NOT IN ('junk') AND ts.name NOT IN ('merged') 
AND q.name NOT IN ('SALES','PRESALES','ODOOHELPDESK') 
AND t.create_time BETWEEN DATE_SUB(NOW(), INTERVAL 31 day) AND NOW() 
AND tt.id IN (2,14)
GROUP BY s.name'''
import numpy as np

df = pd.read_sql(tp_query, con=conn)
df1 = df.T
df1 = df1.reset_index()
df1 = df1.loc[df1.notnull().all(axis=1)]
print(df1)

tp_query2 = '''SELECT  s.name as 'SERVICE REQUEST',
SUM(CASE WHEN ts.id IN (1, 4, 6, 7, 8, 11, 13, 15, 16, 14, 17, 18, 2, 3, 10, 12) THEN 1 ELSE 0 END) AS 'TOTAL',
SUM(CASE WHEN ts.id IN (2, 3, 10, 12) THEN 1 ELSE 0 END) AS 'CLOSED',
SUM(CASE WHEN ts.id IN (1, 4, 6, 7, 8, 11, 13, 15, 16, 14, 17, 18) THEN 1 ELSE 0 END) AS 'Work In Process'
FROM ticket t
JOIN ticket_state ts ON t.ticket_state_id = ts.id
JOIN users u ON t.user_id = u.id
JOIN ticket_type tt ON t.type_id = tt.id
JOIN queue q ON t.queue_id = q.id
left JOIN sla s ON s.id = t.sla_id
WHERE  tt.name NOT IN ('junk') AND ts.name NOT IN ('merged') 
AND q.name NOT IN ('SALES','PRESALES','ODOOHELPDESK') 
AND t.create_time BETWEEN DATE_SUB(NOW(), INTERVAL 31 day) AND NOW() 
AND tt.id IN (4)
GROUP BY s.name'''
import numpy as np

df = pd.read_sql(tp_query2, con=conn)
df2 = df.T
df2 = df2.reset_index()
df2 = df2.loc[df2.notnull().all(axis=1)]
print(df1)

#  add slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
body = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(3))
body.text = 'Ticket Category-Priority Wise'
body.text_frame.paragraphs[0].font.size = Inches(0.3)


# Add table to slide
rows, cols = df1.shape
left = Inches(0)
top = Inches(1)
width = Inches(10)
height = Inches(1)
table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table


# Add table 1 to slide
rows, cols = df1.shape
left = Inches(0)
top = Inches(1)
width = Inches(10)
height = Inches(1)
table_shape1 = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
table1 = table_shape1.table

# Populate data rows of table 1
for i in range(rows):
    for j in range(cols):
        table1.cell(i , j).text = str(df1.values[i, j])
# Set font size of table 1
table_font_size = Pt(13)
for cell in table1.iter_cells():
    cell.text_frame.paragraphs[0].font.size = table_font_size
# Add table 2 to slide
rows, cols = df2.shape
left2 = Inches(0)
top2 = top + height + Inches(2)
width2 = Inches(10)
height2 = Inches(1)
table_shape2 = slide.shapes.add_table(rows + 1, cols, left2, top2, width2, height2)
table2 = table_shape2.table
body.text_frame.paragraphs[0].font.size = Inches(0.5)

# Populate data rows of table 2
for i in range(rows):
    for j in range(cols):
        table2.cell(i,j).text = str(df2.values[i, j])
# Set font size of table 2
for cell in table2.iter_cells():
    cell.text_frame.paragraphs[0].font.size = table_font_size

# end of 4th slide
# start of 5th slide

breached_query = """
SELECT  u.login,s.name,
TIMESTAMPDIFF(MINUTE, t.create_time, t.change_time) AS time_diff,s.solution_time,
CASE
WHEN TIMESTAMPDIFF(MINUTE, t.create_time, t.change_time) > s.solution_time THEN 'Breached'
ELSE 'In time'
END AS status
FROM ticket t
JOIN ticket_state ts ON t.ticket_state_id = ts.id
JOIN users u ON t.user_id = u.id
JOIN ticket_type tt ON t.type_id = tt.id
JOIN queue q ON t.queue_id = q.id
left JOIN sla s ON s.id = t.sla_id
WHERE  tt.name NOT IN ('junk') AND ts.name NOT IN ('merged') 
AND q.name NOT IN ('SALES','PRESALES','ODOOHELPDESK') 
AND t.create_time between DATE_SUB(DATE_SUB(LAST_DAY(CURRENT_DATE), INTERVAL DAY(LAST_DAY(CURRENT_DATE)) - 1 DAY), INTERVAL 1 MONTH) 
and DATE_SUB(DATE_SUB(LAST_DAY(CURRENT_DATE), INTERVAL DAY(LAST_DAY(CURRENT_DATE)) - 1 DAY), INTERVAL 1 DAY) 
AND tt.id IN (14) AND s.id NOT IN(10,11)

"""
cursor.execute(breached_query)
breached_result = cursor.fetchall()

intime_query = """
 SELECT  u.login,s.name,
TIMESTAMPDIFF(MINUTE, t.create_time, t.change_time) AS time_diff,s.solution_time,
CASE
    WHEN TIMESTAMPDIFF(MINUTE, t.create_time, t.change_time) > s.solution_time THEN 'Breached'
    ELSE 'In time'
  END AS status
FROM ticket t
JOIN ticket_state ts ON t.ticket_state_id = ts.id
JOIN users u ON t.user_id = u.id
JOIN ticket_type tt ON t.type_id = tt.id
JOIN queue q ON t.queue_id = q.id
left JOIN sla s ON s.id = t.sla_id
WHERE  tt.name NOT IN ('junk') AND ts.name NOT IN ('merged') 
AND q.name NOT IN ('SALES','PRESALES','ODOOHELPDESK') 
AND t.create_time between DATE_SUB(DATE_SUB(LAST_DAY(CURRENT_DATE), INTERVAL DAY(LAST_DAY(CURRENT_DATE)) - 1 DAY), INTERVAL 1 MONTH) 
and DATE_SUB(DATE_SUB(LAST_DAY(CURRENT_DATE), INTERVAL DAY(LAST_DAY(CURRENT_DATE)) - 1 DAY), INTERVAL 1 DAY) 
AND tt.id IN (4) AND s.id NOT IN(10,11)
"""
cursor.execute(intime_query)
intime_result = cursor.fetchall()
# Set title for slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
body = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(5))
body.text = 'SLA Statistics'
body.text_frame.paragraphs[0].font.size = Inches(0.3)

left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(1)

# Add the first table shape to the slide
table_shape1 = slide.shapes.add_table(rows=2, cols=6, left=left, top=top, width=width, height=height)
table1 = table_shape1.table

# Define the position of the second table shape, which should be below the first table
top2 = top + height + Inches(2)

# Add the second table shape to the slide
table_shape2 = slide.shapes.add_table(rows=2, cols=6, left=left, top=top2, width=width, height=height)
table2 = table_shape2.table

# initialize counters for each status
total = 0
breached_count = 0
in_time_count = 0

# iterate over the list and update the counts
for i in breached_result:

    if i[4] == 'Breached':
        breached_count += 1
    elif i[4] == 'In time':
        in_time_count += 1
    total = (breached_count + in_time_count)
breached_avg = ((breached_count / total) * 100) if total != 0 else 0
breached_avg = round(breached_avg, 2)

# initialize counters for each status
total1 = 0
breached_count1 = 0
in_time_count1 = 0

# iterate over the list and update the counts
for i in intime_result:
    if i[4] == 'Breached':
        breached_count1 += 1
    elif i[4] == 'In time':
        in_time_count1 += 1
    total1 = (breached_count1 + in_time_count1)
    breached_avg1 = (breached_count1 / total1 * 100) if total != 0 else 0
    breached_avg1 = round(breached_avg1, 2)

print(total1)
print(breached_count1)
print(in_time_count1)

# Populate the header row of the table
header_row = table1.rows[0].cells

header_row[0].text = "NUMBER OF INCIDENT "
header_row[1].text = "NUMBER OF INCIDENT WITH SLA ADHERENCE"
header_row[2].text = "NUMBER OF INCIDENT BREACHED SLA"
header_row[3].text = "NUMBER OF SLA WAIVER"
header_row[4].text = "% OF SLA ADHERENCE "
header_row[5].text = "REMARKS"
# Populate the header row of the table
header_row = table2.rows[0].cells

header_row[0].text = "NUMBER OF INCIDENT "
header_row[1].text = "NUMBER OF INCIDENT WITH SLA ADHERENCE"
header_row[2].text = "NUMBER OF INCIDENT BREACHED SLA"
header_row[3].text = "NUMBER OF SLA WAIVER"
header_row[4].text = "% OF SLA ADHERENCE "
header_row[5].text = "REMARKS"

# ...


# Populate the remaining rows of the table with the retrieved data
table_row = table1.rows[1].cells
table_row[0].text = str(total)
table_row[1].text = str(in_time_count)
table_row[2].text = str(breached_count)
table_row[3].text = "-"
table_row[4].text = str(breached_avg)
table_row[5].text = "-"
table_font_size = Pt(13)
# to set size
for cell in table1.iter_cells():
    cell.text_frame.paragraphs[0].font.size = table_font_size


table_row = table2.rows[1].cells
table_row[0].text = str(total1)
table_row[1].text = str(in_time_count1)
table_row[2].text = str(breached_count1)
table_row[3].text = "-"
table_row[4].text = str(breached_avg1)
table_row[5].text = "-"
# to set size
for cell in table2.iter_cells():
    cell.text_frame.paragraphs[0].font.size = table_font_size


# end of 5th slide

# 6th slide
# new_slide_layout = prs.slide_layouts[6]  # Use a different slide layout for variety
# new_slide = prs.slides.add_slide(new_slide_layout)

month1_query = """
 SELECT tt.name,
SUM(CASE WHEN tt.id IN (2,4,14) and t.change_time BETWEEN DATE_SUB(DATE_SUB(LAST_DAY(DATE_SUB(CURRENT_DATE, INTERVAL 3 MONTH)), INTERVAL DAY(LAST_DAY(DATE_SUB(CURRENT_DATE, INTERVAL 3 MONTH))) - 1 DAY), INTERVAL 0 DAY)
AND DATE_SUB(LAST_DAY(DATE_SUB(CURRENT_DATE, INTERVAL 3 MONTH)), INTERVAL 0 DAY) THEN 1 ELSE 0 END) AS 3month,
SUM(CASE WHEN tt.id IN (2,4,14) and t.create_time BETWEEN DATE_SUB(DATE_SUB(LAST_DAY(DATE_SUB(CURRENT_DATE, INTERVAL 2 MONTH)), INTERVAL DAY(LAST_DAY(DATE_SUB(CURRENT_DATE, INTERVAL 2 MONTH))) - 1 DAY), INTERVAL 0 DAY)
AND DATE_SUB(LAST_DAY(DATE_SUB(CURRENT_DATE, INTERVAL 2 MONTH)), INTERVAL 0 DAY) THEN 1 ELSE 0 END) AS 2month,
SUM(CASE WHEN tt.id IN (2,4,14) and t.create_time BETWEEN DATE_SUB(DATE_SUB(LAST_DAY(CURRENT_DATE), INTERVAL DAY(LAST_DAY(CURRENT_DATE)) - 1 DAY), INTERVAL 1 MONTH)
and DATE_SUB(DATE_SUB(LAST_DAY(CURRENT_DATE), INTERVAL DAY(LAST_DAY(CURRENT_DATE)) - 1 DAY), INTERVAL 1 DAY)  THEN 1 ELSE 0 END) AS LAST_monty
FROM ticket t
JOIN ticket_state ts ON t.ticket_state_id = ts.id
JOIN users u ON t.user_id = u.id
JOIN ticket_type tt ON t.type_id = tt.id
JOIN queue q ON t.queue_id = q.id
left JOIN sla s ON s.id = t.sla_id
WHERE  tt.name NOT IN ('junk') AND ts.name NOT IN ('merged')
AND q.name NOT IN ('SALES','PRESALES','ODOOHELPDESK')
#AND t.create_time BETWEEN DATE_SUB(NOW(), INTERVAL 31 day) AND NOW()
AND tt.id IN (2,4,14)
GROUP BY tt.name
"""
cursor.execute(month1_query)
month1_result = cursor.fetchall()
slide = prs.slides.add_slide(prs.slide_layouts[6])
body = slide.shapes.add_textbox(Inches(0.2), Inches(0), Inches(9), Inches(3))
body.text = 'Ticket Trends'
body.text_frame.paragraphs[0].font.size = Inches(0.5)

left = Inches(1)
top = Inches(0.7)
width = Inches(8)
height = Inches(1)

# Create two new table shapes on the slide
table_shape1 = slide.shapes.add_table(rows=len(month1_result) + 1, cols=len(cursor.description), left=left, top=top,
                                      width=width, height=height)

table1 = table_shape1.table
# table2 = table_shape2.table


# Populate the header row of the table
header_row = table1.rows[0].cells
# get current date
today = datetime.date.today()

# get last month
last_month = today - datetime.timedelta(days=30)

# get two months ago
two_months_ago = today - datetime.timedelta(days=60)

# get three months ago
three_months_ago = today - datetime.timedelta(days=90)

# get month names
last_month_name = last_month.strftime('%B')
two_months_ago_name = two_months_ago.strftime('%B')
three_months_ago_name = three_months_ago.strftime('%B')

header_row[0].text = "Category"
header_row[1].text = three_months_ago_name
header_row[2].text = two_months_ago_name
header_row[3].text = last_month_name

for i, row in enumerate(month1_result):
    table_row = table1.rows[i + 1].cells
    table_row[0].text = str(row[0])
    table_row[1].text = str(row[1])
    table_row[2].text = str(row[2])
    table_row[3].text = str(row[3])
    # to set size
table_font_size = Pt(15)
for cell in table1.iter_cells():
    cell.text_frame.paragraphs[0].font.size = table_font_size
    # Create the ChartData object and set the categories
    chart_data = ChartData()
    chart_data.categories = [three_months_ago_name, two_months_ago_name, last_month_name]

    # Populate the data series of the chart
    for row in month1_result:
        values = [float(row[1]), float(row[2]), float(row[3])]
        chart_data.add_series(str(row[0]), values)

# Add the chart to the slide
x, y, cx, cy = Inches(1), Inches(3.3), Inches(7), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
chart.has_legend = True
chart.legend.include_in_layout = False

# end of 6th slide
# thank you slide
# add a new slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

logo_path = '/home/ubuntu/Downloads/images.jpeg'
logo = slide.shapes.add_picture(logo_path, Inches(2), Inches(2.5), Inches(5), Inches(3))

# add a header to the slide
header = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
header.text = 'Thank You'
header.text_frame.paragraphs[0].font.size = Inches(0.4)
# set font size and color

# header.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

# Save the presentation
prs.save('ticket_data.pptx')

# Close the database connection
cursor.close()
conn.close()
